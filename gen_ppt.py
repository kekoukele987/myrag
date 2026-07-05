"""生成 MyRAG 项目介绍 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色主题 ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=GRAY, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return txBox

# ══════════════════════════════════════════════════════════════
# Slide 1: 封面
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
    "MyRAG — 从零构建 RAG 系统", font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
    "LangChain × DeepSeek × RAG 技术学习项目全面解析", font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
    "覆盖 7 大模块 · 30+ 代码示例 · 从入门到进阶", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    "2026 · MyRAG Project", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# Slide 2: 目录
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
    "📖 目录", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

toc_items = [
    "01  基础入门 — 大模型调用与 RAG 初探",
    "02  文档 QA 系统 — 从文档到问答的完整链路",
    "03  模型 I/O — Prompt、模型、输出解析",
    "04  提示模板（上）— PromptTemplate 与 Few-Shot",
    "05  提示模板（下）— Chain of Thought 思维链",
    "06  调用模型 — 云端 API 与本地模型",
    "07  解析输出 — 结构化输出解析器",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(5), toc_items, font_size=20, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slide 3: 项目概览
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
    "🔭 项目概览", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

overview = [
    "▸ 目标：系统化学习 RAG（检索增强生成）技术栈",
    "▸ 核心框架：LangChain（1.x）",
    "▸ 大模型：DeepSeek Chat API（兼容 OpenAI 格式）",
    "▸ 向量模型：BAAI/bge-small-zh-v1.5（本地嵌入）",
    "▸ 本地推理：Ollama（Qwen 等模型）",
    "▸ 知识库：本地文档（PDF / Word / TXT）",
    "▸ Web 框架：Flask（文档 QA 系统界面）",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(5), overview, font_size=18, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slides 4-8: 目录 01 基础入门
# ══════════════════════════════════════════════════════════════
# Slide 4: 01 模块标题
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.6),
    "📁 01 / 基础入门", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.8),
    "从最简调用到 RAG 链路", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Slide 5: 01 文件详解
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
    "01 / 基础入门 — 文件详解", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

files_01 = [
    "01_简单文本生成.py    — 最简 LLM 调用（DeepSeek Chat，invoke 单条文本）",
    "02_看图说话.py        — 多模态：BLIP 图片描述 + LLM 文案生成",
    "02_ChatModel.py       — 原生 OpenAI SDK 调用 Chat 模型（messages 格式）",
    "03_TextLangChain_v0.2 — LangChain 链式调用（prompt → model → output）",
    "04_ChatLongChain.py   — SystemMessage + HumanMessage 多角色对话",
    "lang.py               — LangChain 完整封装：BGE 向量化 + Chroma 检索",
    "rag.py                — LlamaIndex RAG：自定义 DeepSeek LLM 包装类",
    "设定.txt              — 《黑神话：悟空》背景知识（测试用文档）",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), files_01, font_size=16, color=GRAY)

# Slide 6: 01 知识点
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
    "01 / 核心知识点", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), GREEN)

points_01 = [
    "✦ LLM 调用的三种方式：",
    "   · 原生 OpenAI SDK — client.chat.completions.create()",
    "   · LangChain ChatOpenAI — llm.invoke(prompt)",
    "   · LlamaIndex — 自定义 LLM 包装类",
    "",
    "✦ 多模态理解：BLIP（图片转描述）+ LLM（描述转文案）",
    "",
    "✦ LangChain 核心模块：",
    "   · HuggingFaceEmbeddings — 本地向量化（BGE）",
    "   · Chroma — 本地向量数据库",
    "   · RecursiveCharacterTextSplitter — 文档切分",
    "   · LCEL — LangChain Expression Language 链式调用",
    "",
    "✦ RAG 核心：检索 + 生成（本地检索 → 云端生成）",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), points_01, font_size=15, color=GRAY)

# Slide 7: 01 运行效果
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
    "01 / 运行效果展示", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

demo_01 = [
    "▸ 01_简单文本生成.py：",
    "    输入 → \"请给我写一句清明节的中文宣传语\"",
    "    输出 → 模型生成的宣传文案",
    "",
    "▸ lang.py（RAG 链路）：",
    "    查询 → \"bge模型和deepseek api分别承担什么作用？\"",
    "    本地检索 → BGE 匹配到最相关文档",
    "    生成回答 → DeepSeek 结合检索结果生成答案",
    "",
    "▸ 02_看图说话.py：",
    "    输入 → 图片 URL（Pexels 免费图库）",
    "    BLIP 生成英文描述 → DeepSeek 生成中文推广文案",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), demo_01, font_size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slides 8-11: 目录 02 文档 QA 系统
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.6),
    "📁 02 / 文档 QA 系统", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.8),
    "RAG 完整链路的 Web 实现", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "02 / 文档 QA 系统 — 架构图", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

arch_items = [
    "RAG 完整链路（DocQA_v0.1.py）：",
    "",
    "  ① Load（加载） → 支持 PDF / DOCX / TXT 格式",
    "  ② Split（切分） → RecursiveCharacterTextSplitter（chunk_size=500）",
    "  ③ Store（存储） → FastEmbed（BGE）→ Chroma 向量数据库",
    "  ④ Retrieval（检索） → 语义相似度匹配 Top-4",
    "  ⑤ Augment（增强） → Prompt 拼接（上下文 + 问题）",
    "  ⑥ Generation（生成） → DeepSeek Chat API",
    "  ⑦ Output（输出） → Web 界面（Flask + HTML 模板）",
    "",
    "知识库：OneFlower/ 目录下的花语大全、员工手册、运营指南",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), arch_items, font_size=16, color=GRAY)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "02 / 核心知识点", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), GREEN)

points_02 = [
    "✦ Document Loaders：PyPDFLoader / Docx2txtLoader / TextLoader",
    "✦ FastEmbedEmbeddings：无需联网下载模型，本地快速嵌入",
    "✦ Chroma 向量数据库：本地持久化存储，支持相似度检索",
    "✦ DeepSeek Chat：作为生成模型，temperature=0 保证确定性",
    "✦ Flask Web 界面：提供用户友好的问答交互",
    "✦ 异常处理：加载失败时跳过，保证系统鲁棒性",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.5), points_02, font_size=17, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slides 12-15: 目录 03 模型 I/O
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.6),
    "📁 03 / 模型 I/O", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.8),
    "Prompt → LLM → Output 的多种形态", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "03 / 模型 I/O — 文件详解", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

files_03 = [
    "01_模型IO.py           — PromptTemplate + ChatOpenAI 单次调用",
    "02_模型IO_循环调用.py   — 循环调用批量生成（玫瑰/百合/康乃馨）",
    "03_OpenAI_IO.py        — 原生 OpenAI SDK 实现（不用 LangChain）",
    "04_模型IO_HuggingFace  — 替换 HuggingFace 为 DeepSeek（同 01 结构）",
    "05_模型IO_输出解析.py   — JsonOutputParser 结构化输出 + CSV 保存",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4), files_03, font_size=17, color=GRAY)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "03 / 四种调用方式对比", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ORANGE)

compare_03 = [
    "┌─────────────────┬──────────────────┬──────────────────┐",
    "│     文件         │   调用方式        │   输出类型        │",
    "├─────────────────┼──────────────────┼──────────────────┤",
    "│ 01_模型IO        │ ChatOpenAI.invoke │ 纯文本 str        │",
    "│ 02_循环调用      │ for 循环 + invoke │ 批量纯文本        │",
    "│ 03_OpenAI_IO     │ OpenAI SDK       │ 纯文本 str        │",
    "│ 05_输出解析      │ + JsonOutputParser│ 结构化 dict       │",
    "└─────────────────┴──────────────────┴──────────────────┘",
    "",
    "✦ 05 引入 JsonOutputParser：要求 LLM 按 JSON 格式返回",
    "✦ 解析后存入 DataFrame → 保存为 CSV（flowers_with_descriptions.csv）",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), compare_03, font_size=16, color=GRAY)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "03 / 核心知识点", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), GREEN)

points_03 = [
    "✦ PromptTemplate：模板字符串 + format() 填充变量",
    "✦ ChatOpenAI vs OpenAI SDK：封装层 vs 原生层",
    "✦ 循环调用：批量场景下的模型调用模式",
    "✦ JsonOutputParser：prompt 中注入 JSON Schema 约束",
    "✦ 输出持久化：结构化数据 → DataFrame → CSV",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.5), points_03, font_size=17, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slides 16-19: 目录 04 提示模板上
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.6),
    "📁 04 / 提示模板（上）", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.8),
    "Prompt 工程的三种基础模板", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "04 / 提示模板（上）— 文件详解", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

files_04 = [
    "00_ImportTemplates.py     — 演示各种 PromptTemplate 的导入方式",
    "01_PromptTemplate.py      — 两种创建方式：from_template / 直接构造",
    "02_ChatPromptTemplate.py  — SystemMessage + HumanMessage 多角色模板",
    "03_FewShotPrompt.py       — 少样本提示 + 语义相似度示例选择器",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3), files_04, font_size=17, color=GRAY)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "04 / FewShotPrompt 详解", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

fewshot = [
    "✦ 示例数据（samples）：玫瑰→爱情、康乃馨→母亲节、百合→庆祝、向日葵→鼓励",
    "",
    "✦ 两种提供示例的方式：",
    "   · 方式一：直接指定 examples= 列表，所有示例拼入 prompt",
    "   · 方式二：SemanticSimilarityExampleSelector",
    "     基于语义相似度自动选择最相关的 K 个示例",
    "     使用 BGE 向量模型 + Chroma 向量库",
    "",
    "✦ 效果：输入\"野玫瑰\" → 自动匹配\"玫瑰\"示例，生成更精准的文案",
    "",
    "✦ LangChain 版本迁移：",
    "   · langchain.prompts → langchain_core.prompts",
    "   · OpenAIEmbeddings → FastEmbedEmbeddings（本地 BGE）",
    "   · Qdrant → Chroma（轻量本地向量库）",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), fewshot, font_size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slides 20-22: 目录 05 提示模板下
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.6),
    "📁 05 / 提示模板（下）", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.8),
    "Chain of Thought 思维链推理", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "05 / Chain of Thought 详解", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

cot = [
    "✦ CoT（思维链）核心思想：",
    "   让 LLM 在给出最终答案之前，先展示推理过程",
    "   模仿人类\"一步步思考\"的方式，提升推理准确性",
    "",
    "✦ 本示例的 CoT 模板结构：",
    "   · SystemMessage（角色设定）：花店电商 AI 助手",
    "   · SystemMessage（CoT 规则）：",
    "     先理解需求 → 考虑鲜花涵义 → 给出推荐 → 解释原因",
    "   · HumanMessage：用户输入（如\"女朋友喜欢粉色和紫色\"）",
    "",
    "✦ 结合 Few-Shot 学习：",
    "   CoT 模板中嵌入两个示例（象征爱情→红玫瑰，独特花朵→兰花）",
    "   每个示例展示完整的推理步骤",
    "",
    "✦ 参数设置：temperature=0（确定性输出，适合推理场景）",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), cot, font_size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slides 23-25: 目录 06 调用模型
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.6),
    "📁 06 / 调用模型", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.8),
    "云端 API 与本地模型多种调用方式", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "06 / 文件详解与对比", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

files_06 = [
    "01_HugggingFace_Llama.py       — 原调用 Llama-2 本地模型，现改为 DeepSeek API",
    "02_LangChain_HFHub.py          — LCEL 链式调用（prompt | llm | parser）",
    "04_LangChain_CustomizeModel.py — 两种调用本地模型的方式",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2), files_06, font_size=17, color=GRAY)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "06 / 自定义模型调用详解", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), GREEN)

custom = [
    "✦ 方式一：ChatOllama（推荐）",
    "   直接使用 langchain_ollama.ChatOllama",
    "   配置 model_name / temperature / base_url",
    "   适合快速集成 Ollama 管理的本地模型",
    "",
    "✦ 方式二：自定义 CustomLLM 类",
    "   继承 langchain_core.language_models.llms.LLM",
    "   实现 _call() 抽象方法 → 统一 invoke() 接口",
    "   适合需要替换 LLM 后端的生产场景",
    "",
    "✦ 重构对比（llama_cpp → Ollama）：",
    "   · 模型加载：硬编码 GGUF 文件 → Ollama 管理模型生命周期",
    "   · 平台兼容：Linux only → 跨平台（Win/Mac/Linux）",
    "   · 模型切换：改路径 → 改一行模型名",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), custom, font_size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slides 26-28: 目录 07 解析输出
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(5), Inches(2.8), Inches(3.3), Inches(0.6),
    "📁 07 / 解析输出", font_size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5), Inches(3.6), Inches(3.3), Inches(0.8),
    "从非结构化文本到结构化数据", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "07 / 文件详解与对比", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

files_07 = [
    "01_Pydantic_Parser.py   — PydanticOutputParser：基于 Pydantic 模型的 JSON 解析",
    "02_OutputFixParser.py    — OutputFixingParser：用 LLM 自动修复格式错误",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.5), files_07, font_size=17, color=GRAY)

# 新增对比 slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "07 / 三种输出解析器对比", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ORANGE)

compare_07 = [
    "┌────────────────────┬──────────────────┬──────────────────────┐",
    "│     解析器          │   输入格式        │   特性               │",
    "├────────────────────┼──────────────────┼──────────────────────┤",
    "│ JsonOutputParser   │ 任意文本          │ 提取 JSON → dict    │",
    "│ PydanticParser     │ 合法 JSON         │ 类型校验 + 字段描述  │",
    "│ OutputFixingParser │ 不合法 JSON       │ 自动修复 + 类型校验  │",
    "└────────────────────┴──────────────────┴──────────────────────┘",
    "",
    "✦ OutputFixingParser 工作原理：",
    "  ① 先用 PydanticOutputParser 尝试解析",
    "  ② 解析失败 → 将错误信息 + 错误输出发给 LLM 修正",
    "  ③ LLM 返回修正版 → 重新解析",
    "  ④ 本质：用一次额外 LLM 调用换取解析健壮性",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5), compare_07, font_size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slide 29: 技术栈总结
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "🛠️ 技术栈总结", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT)

tech_stack = [
    "✦ LangChain 1.x — 核心框架（prompt / model / chain / output_parser）",
    "✦ DeepSeek Chat API — 云端大模型（兼容 OpenAI 格式）",
    "✦ BAAI/bge-small-zh-v1.5 — 中文语义向量模型（本地嵌入）",
    "✦ Chroma — 本地向量数据库（持久化存储）",
    "✦ Ollama — 本地模型运行框架（Qwen / Llama 等）",
    "✦ FastEmbed — 轻量嵌入库（无需手动下载模型文件）",
    "✦ BLIP — 图片描述生成（多模态理解）",
    "✦ Flask — Web 服务框架（文档 QA 系统界面）",
    "✦ Pydantic — 数据模型定义与校验（结构化输出）",
    "✦ LlamaIndex — 替代 RAG 框架（rag.py 中对比使用）",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), tech_stack, font_size=17, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slide 30: 学习路线图
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
    "🗺️ 推荐学习路线", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), GREEN)

roadmap = [
    "第一阶段：基础入门（01/）",
    "   大模型调用基础 → Chat 模型 → LangChain 链式调用 → RAG 初体验",
    "",
    "第二阶段：文档 QA 系统（02/）",
    "   完整 RAG 链路：加载 → 切分 → 向量化 → 存储 → 检索 → 生成 → 展示",
    "",
    "第三阶段：模型 I/O（03/）",
    "   PromptTemplate → 循环调用 → OpenAI SDK → 输出解析 → 持久化",
    "",
    "第四阶段：提示工程（04/ 05/）",
    "   PromptTemplate → ChatPromptTemplate → Few-Shot → CoT 思维链",
    "",
    "第五阶段：模型调用进阶（06/）",
    "   云端 API → LCEL 链式调用 → 本地模型 → 自定义 LLM 封装",
    "",
    "第六阶段：输出解析（07/）",
    "   JsonOutputParser → PydanticOutputParser → OutputFixingParser",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.5), roadmap, font_size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════
# Slide 31: 结束页
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), ACCENT)
add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1),
    "Thank You！", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
    "MyRAG — 从零构建 RAG 系统学习项目", font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
    "https://github.com/kekoukele987/myrag", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = "MyRAG_项目介绍.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")